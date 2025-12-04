from __future__ import annotations
# -*- coding: utf-8 -*-

"""
word_to_pptx_treinamento.py (VERSÃO 6.0 - Preparado para VBA)

- Cria a apresentação com o tema Salvioli (sem faixa azul).
- IMPORTANTE: O VBA 'ReformatarTituloEBullets' DEVE ser executado
  após a criação do PPTX para ajustar a posição e o tamanho dinâmico das fontes.
- Utiliza o LAYOUT 1 (Título e Conteúdo) sem forçar o redimensionamento dos placeholders.
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


# =========================
# CONFIGURAÇÕES DE TEMA E LAYOUT
# =========================
THEME_COLOR = RGBColor(0, 51, 102)          # Azul Salvioli
BACKGROUND_COLOR = RGBColor(242, 242, 242)  # Cinza claro
FONT_NAME = "Arial"

MAX_BULLETS = 15  # Máximo de bullets por slide
LOGO_SIZE = Inches(1.0)     


# =========================
# FUNÇÕES DE TEMA
# =========================
def apply_slide_background(slide) -> None:
    """
    Aplica cor de fundo no slide.
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR


def add_top_banner(slide, slide_width) -> None:
    """
    Faixa superior azul removida a pedido do usuário.
    """
    pass


def add_logo(slide, logo_path: Optional[Path], slide_width, slide_height) -> None:
    """
    Adiciona o logo no canto inferior direito do slide.
    """
    if not logo_path or not logo_path.exists():
        return

    # Posição: Canto inferior direito, com margem de 0.3 polegadas (0.76 cm)
    left = slide_width - LOGO_SIZE - Inches(0.3)
    top = slide_height - LOGO_SIZE - Inches(0.3)

    slide.shapes.add_picture(str(logo_path), left, top, width=LOGO_SIZE, height=LOGO_SIZE)


def add_formatted_title(body_shape, text: str) -> None:
    """
    Formata o texto do Title Placeholder (índice 0) para o tema Salvioli.
    """
    if body_shape is None:
        return

    # O VBA ajustará a posição e o tamanho do placeholder, o Python apenas formata o texto
    tf = body_shape.text_frame
    tf.text = text
    
    tf.word_wrap = True
    # O VBA sobrescreverá o AutoSize, mas definimos o padrão aqui.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Título será centralizado

    p = tf.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(28) 
    p.font.bold = True
    p.font.color.rgb = THEME_COLOR 
    p.alignment = PP_ALIGN.CENTER # O VBA muda isso para CENTER, vamos usar o padrão CENTER


def add_bullet_to_body(body_shape, text: str, style_level: int = 0, is_highlight: bool = False) -> None:
    """
    Adiciona um bullet (parágrafo) no corpo do slide.
    """
    tf = body_shape.text_frame
    
    if not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.text = text
    p.level = style_level
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT

    if is_highlight:
        # Formatação para Heading 2 (Destaque)
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = THEME_COLOR
    else:
        # Formatação para Texto Normal (Bullet) - VBA ajustará para caber
        p.font.size = Pt(20) 
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0)


def add_slide_numbers(pres: Presentation) -> None:
    """
    Adiciona numeração de slides ("n / total") no canto inferior esquerdo.
    """
    total = len(pres.slides)
    if total == 0:
        return

    width = pres.slide_width
    height = pres.slide_height

    for idx, slide in enumerate(pres.slides, start=1):
        tx_width = Inches(1.8)
        tx_height = Inches(0.5)
        
        # Posição: no rodapé, no canto esquerdo
        left = Inches(0.3)
        top = height - tx_height - Inches(0.1) 

        textbox = slide.shapes.add_textbox(left, top, tx_width, tx_height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{idx} / {total}"
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.alignment = PP_ALIGN.LEFT


# =========================
# CORE: CONVERSÃO DOCX -> PPTX
# =========================
def _novo_slide_conteudo(pres: Presentation,
                         logo_path: Optional[Path],
                         slide_width,
                         slide_height,
                         titulo: str) -> tuple:
    """
    Cria um novo slide de conteúdo, utilizando o LAYOUT 1 (Título e Conteúdo)
    """
    slide_layout = pres.slide_layouts[1]  # Layout 1: Título e Conteúdo
    slide = pres.slides.add_slide(slide_layout)

    # 1. Aplica Tema (Fundo, Logo)
    apply_slide_background(slide)
    add_top_banner(slide, slide_width) 
    add_logo(slide, logo_path, slide_width, slide_height)
    
    # 2. Obtém os Placeholders padrão que o VBA irá ajustar
    title_shape = None
    body_shape = None
    try:
        # Placeholder de Título (índice 0)
        title_shape = slide.placeholders[0]
        # Placeholder de Conteúdo (índice 1)
        body_shape = slide.placeholders[1]
    except IndexError:
        print("Aviso: Placeholders padrão (0 ou 1) não encontrados.")
        
    
    # 3. Formata o Título
    add_formatted_title(title_shape, titulo)
    
    # 4. Configuração inicial de texto para o VBA trabalhar
    if body_shape:
        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE # Desabilita auto-ajuste para o VBA ter controle

    return slide, body_shape


def create_presentation_from_docx(
    docx_path: Path,
    pptx_path: Path,
    logo_path: Optional[Path] = None,
    titulo_padrao: Optional[str] = None,
) -> None:
    """
    Converte um DOCX em PPTX.
    """
    document = Document(str(docx_path))
    pres = Presentation()

    slide_width = pres.slide_width
    slide_height = pres.slide_height

    current_slide = None
    body_shape = None
    bullet_count = 0
    titulo_modulo = None

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""
        
        # Cria slide inicial se ainda não existir
        if current_slide is None:
            titulo_modulo = titulo_padrao or docx_path.stem.replace("_", " ").title()
            current_slide, body_shape = _novo_slide_conteudo(
                pres, logo_path, slide_width, slide_height, titulo_modulo
            )
            if body_shape is None:
                continue

        # Heading 1 => novo módulo (novo slide principal)
        if style_name.startswith("Heading 1") or style_name.startswith("Título 1"):
            titulo_modulo = text
            bullet_count = 0
            current_slide, body_shape = _novo_slide_conteudo(
                pres, logo_path, slide_width, slide_height, titulo_modulo
            )

        if body_shape:
            # Heading 2 => subtítulo em destaque dentro do módulo
            if style_name.startswith("Heading 2") or style_name.startswith("Título 2"):
                # Quebra em continuação se MAX_BULLETS atingido
                if bullet_count >= MAX_BULLETS:
                    titulo_cont = (titulo_modulo or titulo_padrao or docx_path.stem) + " (continuação)"
                    bullet_count = 0
                    current_slide, body_shape = _novo_slide_conteudo(
                        pres, logo_path, slide_width, slide_height, titulo_cont
                    )
                
                if body_shape:
                    add_bullet_to_body(body_shape, text, is_highlight=True)
                    bullet_count += 1

            # Texto normal => bullets
            else:
                # Quebra em continuação se MAX_BULLETS atingido
                if bullet_count >= MAX_BULLETS:
                    titulo_cont = (titulo_modulo or titulo_padrao or docx_path.stem) + " (continuação)"
                    bullet_count = 0
                    current_slide, body_shape = _novo_slide_conteudo(
                        pres, logo_path, slide_width, slide_height, titulo_cont
                    )
                
                if body_shape:
                    add_bullet_to_body(body_shape, text)
                    bullet_count += 1

    # Slide final se o DOCX estava vazio
    if len(pres.slides) == 0:
        titulo = titulo_padrao or docx_path.stem.replace("_", " ").title()
        _novo_slide_conteudo(
            pres, logo_path, slide_width, slide_height, titulo
        )

    # Numeração
    add_slide_numbers(pres)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(pptx_path))
    print(f"[OK] {docx_path.name} -> {pptx_path}")

# =========================
# PROCESSAMENTO DE PASTA / MAIN
# =========================
# (Mantido o código original)

def processar_pasta(pasta_base: Path, logo: Optional[Path], titulo: Optional[str]) -> None:
    pasta_base = pasta_base.resolve()
    out_base = pasta_base.parent / "gerados_pptx"
    out_base.mkdir(parents=True, exist_ok=True)

    arquivos = list(pasta_base.rglob("*.docx"))
    if not arquivos:
        print("Nenhum arquivo .docx encontrado na pasta (incluindo subpastas).")
        return

    print(f"Encontrados {len(arquivos)} arquivos .docx em {pasta_base} (com subpastas).")

    for docx_file in arquivos:
        rel = docx_file.relative_to(pasta_base)
        rel_pptx = rel.with_suffix(".pptx")
        pptx_destino = out_base / rel_pptx
        pptx_destino.parent.mkdir(parents=True, exist_ok=True)

        titulo_padrao = titulo or docx_file.stem.replace("_", " ").title()

        create_presentation_from_docx(
            docx_path=docx_file,
            pptx_path=pptx_destino,
            logo_path=logo,
            titulo_padrao=titulo_padrao
        )

    print(f"Apresentações geradas em: {out_base}")


def main():
    parser = argparse.ArgumentParser(
        description="Converte DOCX (arquivo único ou pasta) em PPTX com tema Salvioli."
    )

    parser.add_argument("--docx", help="Arquivo DOCX único.")
    parser.add_argument("--pptx", help="Arquivo PPTX de saída.")
    parser.add_argument("--pasta", help="Pasta base contendo arquivos .docx (processa recursivamente).")
    parser.add_argument("--logo", help="Logo da empresa (PNG/JPG).")
    parser.add_argument("--titulo", help="Título padrão (usado quando não houver Heading 1).")

    args = parser.parse_args()

    logo = Path(args.logo).resolve() if args.logo else None
    if logo and not logo.exists():
        print(f"Aviso: logo não encontrado em {logo}. Continuando sem logo.")
        logo = None

    if args.pasta:
        pasta = Path(args.pasta).resolve()
        if not pasta.exists():
            print(f"Erro: pasta não encontrada: {pasta}")
            return
        processar_pasta(pasta, logo, args.titulo)
        return

    if args.docx and args.pptx:
        docx = Path(args.docx).resolve()
        if not docx.exists():
            print(f"Erro: DOCX não encontrado: {docx}")
            return
        pptx = Path(args.pptx).resolve()

        create_presentation_from_docx(
            docx_path=docx,
            pptx_path=pptx,
            logo_path=logo,
            titulo_padrao=args.titulo
        )
        return

    if len(sys.argv) == 1:
        script_dir = Path(__file__).resolve().parent
        pasta_auto = script_dir / "MODULOS"
        if pasta_auto.exists():
            print(f"Nenhum argumento informado. Processando pasta padrão: {pasta_auto}")
            processar_pasta(pasta_auto, logo, args.titulo)
        else:
            print(
                "Nenhum argumento informado e pasta padrão 'MODULOS' não encontrada.\n"
                "Use um dos modos:\n"
                "  python word_to_pptx_treinamento.py --pasta CAMINHO_DOS_DOCX\n"
                "  python word_to_pptx_treinamento.py --docx ARQ.docx --pptx ARQ.pptx"
            )
        return

    print("Parâmetros insuficientes. Informe --pasta ou (--docx e --pptx).")


if __name__ == "__main__":
    main()